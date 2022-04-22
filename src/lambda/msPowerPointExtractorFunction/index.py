import json
import urllib.parse
import os
from pptx import Presentation
from common import *

print('Loading function')


def extract_text(file_path):
    prs = Presentation(file_path)
    texts = {}
    slideNumber = 1
    for slide in prs.slides:
        slideTexts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
                slideTexts.append(shape.text)
        texts[slideNumber] = "/n".join(slideTexts)
        slideNumber += 1
    print(texts)
    filename, file_extension = os.path.splitext(file_path)
    output_file = os.path.join(os.path.dirname(file_path), filename + ".json")
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(texts, f)


def lambda_handler(event, context):
    clean_tmp()
    # Get the object from the event and show its content type
    bucket = event['Records'][0]['s3']['bucket']['name']
    key = urllib.parse.unquote_plus(
        event['Records'][0]['s3']['object']['key'], encoding='utf-8')
    try:
        file_path = save_file(bucket, key)
        extract_media(file_path, 'ppt/media/')
        extract_text(file_path)
        copy_tmp_to_processing_bucket()
        return 'OK'
    except Exception as e:
        print(e)
        print('Error getting object {} from bucket {}. Make sure they exist and your bucket is in the same region as this function.'.format(key, bucket))
        raise e
